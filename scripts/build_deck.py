# -*- coding: utf-8 -*-
"""build_deck.py — 按 plan.json 克隆模板页并填充内容, 输出成品 PPTX。

用法:
    # 1) 生成计划骨架(每个非装饰文本槽预填「【待填】+ 模板原文」, 填空即可)
    python build_deck.py <plan.json> --scaffold --template <模板.pptx> --pages 1,2,4,4,6

    # 2) 校验 / 生成
    python build_deck.py <plan.json> [--check] [--allow-residual]

plan.json 结构(详见 references/plan-format.md):
{
  "template": "模板.pptx 路径",
  "output":   "成品.pptx 路径",
  "slides": [
    {
      "source": 3,                          // 模板页码(1-based), 可重复使用
      "texts":  {"14": "新标题\\n副标题",    // 文本槽: shape_id -> 新文本(\\n 分段)
                 "9!r2c3": "单元格内容"},    // 表格槽: shape_id!r行c列
      "images": {"22": "assets/foo.png"},   // 图片槽: shape_id -> 图片路径(相对 plan.json)
      "notes":  "演讲者备注(可选)"
    }
  ]
}

--check 只做校验不生成。校验错误(E)会阻止生成; 警告(W)不阻止但必须复核。
"""

import argparse
import io
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

from pptx import Presentation  # noqa: E402

from pptxkit import (  # noqa: E402
    collect_slots, duplicate_slide, delete_slide, find_shape,
    set_shape_text, set_table_cell_text, replace_picture,
)

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"}

SENTINEL = "【待填】"


def scaffold(plan_path, template, pages, deck_out):
    """按选定页序生成 plan.json 骨架: 非装饰文本槽预填「【待填】+ 原文」。"""
    prs = Presentation(str(template))
    total = len(prs.slides)
    slides_json = []
    for src in pages:
        if not (1 <= src <= total):
            print(f"错误: --pages 中的 {src} 不在模板页码范围 1-{total}", file=sys.stderr)
            sys.exit(2)
        slots = collect_slots(prs.slides[src - 1])
        texts, hints = {}, {}
        for t in slots["text_slots"]:
            if t["decorative"]:
                continue
            sample = t["sample"].strip()
            texts[t["id"]] = f"{SENTINEL} {t['sample']}" if sample else SENTINEL
            h = [t["name"] or "?", f"样本{t['sample_len']}字"]
            if t.get("max_chars"):
                h.append(f"预算≤{t['max_chars']}字")
            if t.get("max_lines"):
                h.append(f"≤{t['max_lines']}行")
            if t.get("paragraphs", 1) > 1:
                h.append(f"{t['paragraphs']}段")
            if t.get("layout_prompt"):
                h.append(f"用途:{t['layout_prompt']}")
            hints[t["id"]] = " | ".join(h)
        for tb in slots["table_slots"]:
            for cell_key, cell_sample in tb["cells"].items():
                texts[f"{tb['id']}!{cell_key}"] = cell_sample
            hints[tb["id"]] = (f"表格{tb['rows']}x{tb['cols']}: 各格已按原文预填, "
                               f"只改需要变的格子")
        entry = {"source": src}
        if hints:
            entry["_hints"] = hints
        if slots["image_slots"]:
            entry["_image_slots"] = {
                i["id"]: f"{i['name']} 宽高比{i['aspect']}" for i in slots["image_slots"]
            }
        if texts:
            entry["texts"] = texts
        entry["notes"] = ""
        slides_json.append(entry)
    plan = {
        "template": str(Path(template).resolve()),
        "output": str(deck_out) if deck_out else "output.pptx",
        "slides": slides_json,
    }
    plan_path.parent.mkdir(parents=True, exist_ok=True)
    plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2),
                         encoding="utf-8")
    n_todo = sum(v.count(SENTINEL) for s in slides_json
                 for v in (s.get("texts") or {}).values())
    print(f"骨架已写入 {plan_path}: {len(slides_json)} 页, {n_todo} 个{SENTINEL}槽位。\n"
          f"下一步: 把每个「{SENTINEL}」值整体替换为新文案"
          f"(保留模板原文 = 只删掉「{SENTINEL} 」前缀), 然后 --check。")


def parse_slot_id(slot_id):
    """返回 (shape_id:int, row:int|None, col:int|None)。"""
    s = str(slot_id)
    if "!" in s:
        sid, cell = s.split("!", 1)
        import re
        m = re.fullmatch(r"r(\d+)c(\d+)", cell)
        if not m:
            raise ValueError(f"非法表格槽位 ID: {slot_id} (应为 <shape_id>!r<行>c<列>)")
        return int(sid), int(m.group(1)), int(m.group(2))
    return int(s), None, None


def validate(plan, plan_dir, prs, allow_residual=False):
    errors, warnings = [], []
    slides = plan.get("slides")
    if not isinstance(slides, list) or not slides:
        return ["plan.slides 必须是非空数组"], warnings
    total = len(prs.slides)
    slot_cache = {}

    for pi, entry in enumerate(slides, 1):
        src = entry.get("source")
        if not isinstance(src, int) or not (1 <= src <= total):
            errors.append(f"成品第{pi}页: source={src!r} 不在模板页码范围 1-{total}")
            continue
        if src not in slot_cache:
            slot_cache[src] = collect_slots(prs.slides[src - 1])
        slots = slot_cache[src]
        text_by_id = {t["id"]: t for t in slots["text_slots"]}
        table_by_id = {t["id"]: t for t in slots["table_slots"]}
        image_by_id = {t["id"]: t for t in slots["image_slots"]}

        filled_ids = set()
        for slot_id, text in (entry.get("texts") or {}).items():
            try:
                sid, row, col = parse_slot_id(slot_id)
            except ValueError as e:
                errors.append(f"成品第{pi}页: {e}")
                continue
            if SENTINEL in str(text):
                errors.append(
                    f"成品第{pi}页 槽{slot_id}: 内容仍含「{SENTINEL}」, 该槽尚未填写。"
                    f"整体替换为新文案, 或删掉前缀以保留模板原文")
            if row is not None:
                tbl = table_by_id.get(str(sid))
                if tbl is None:
                    errors.append(f"成品第{pi}页: 模板第{src}页没有表格槽 {sid}")
                elif not (1 <= row <= tbl["rows"] and 1 <= col <= tbl["cols"]):
                    errors.append(
                        f"成品第{pi}页: 表格{sid}尺寸 {tbl['rows']}x{tbl['cols']}, "
                        f"r{row}c{col} 越界")
                continue
            slot = text_by_id.get(str(sid))
            if slot is None:
                errors.append(f"成品第{pi}页: 模板第{src}页没有文本槽 {sid}"
                              f"(可用: {', '.join(text_by_id) or '无'})")
                continue
            if slot.get("decorative"):
                warnings.append(
                    f"成品第{pi}页 槽{sid}({slot['name']}): 这是装饰槽"
                    f"(页码/序号/LOGO), 不应填写, 建议从 plan 中删除")
            filled_ids.add(str(sid))
            plain_len = len(str(text).replace("\n", ""))
            budget = slot.get("max_chars")
            if budget is None and slot["sample_len"]:
                budget = round(slot["sample_len"] * 1.3)
            if budget and plain_len > budget * 1.15:
                warnings.append(
                    f"成品第{pi}页 槽{sid}({slot['name']}): 文本 {plain_len} 字超出预算 "
                    f"{budget}(样本 {slot['sample_len']} 字), 可能溢出/换行难看")
            if slot.get("paragraphs", 1) == 1 and slot.get("max_lines") == 1 \
                    and "\n" in str(text):
                warnings.append(f"成品第{pi}页 槽{sid}: 单行槽填入了多段文本")

        for slot_id, img in (entry.get("images") or {}).items():
            slot = image_by_id.get(str(slot_id))
            if slot is None:
                errors.append(f"成品第{pi}页: 模板第{src}页没有图片槽 {slot_id}"
                              f"(可用: {', '.join(image_by_id) or '无'})")
                continue
            p = (plan_dir / img).resolve() if not Path(img).is_absolute() else Path(img)
            if not p.is_file():
                errors.append(f"成品第{pi}页 图片槽{slot_id}: 找不到图片 {p}")
            elif p.suffix.lower() not in IMAGE_EXTS:
                warnings.append(f"成品第{pi}页 图片槽{slot_id}: 非常见图片格式 {p.suffix}")

        # 未填的非装饰文本槽 = 模板演示文案会残留。默认按错误阻断:
        # 弱执行环境里警告常被无视, 而文案残留是最常见、观感最差的交付事故。
        unfilled = [t for t in slots["text_slots"]
                    if not t["decorative"] and t["id"] not in filled_ids
                    and t["sample"].strip()]
        for t in unfilled:
            msg = (f"成品第{pi}页(模板第{src}页) 槽{t['id']}({t['name']}) 未填, "
                   f"模板文案将残留: “{t['sample'][:20]}…”。"
                   f"保留原文请把原文填入该槽")
            if allow_residual:
                warnings.append(msg)
            else:
                errors.append(msg + "; 确认接受残留可加 --allow-residual")
    return errors, warnings


def build(plan, plan_dir, prs):
    n_template = len(prs.slides)
    for entry in plan["slides"]:
        src_slide = prs.slides[entry["source"] - 1]
        new_slide = duplicate_slide(prs, src_slide)
        for slot_id, text in (entry.get("texts") or {}).items():
            sid, row, col = parse_slot_id(slot_id)
            shape = find_shape(new_slide, sid)
            if row is not None:
                set_table_cell_text(shape, row, col, text)
            else:
                set_shape_text(shape, text)
        for slot_id, img in (entry.get("images") or {}).items():
            shape = find_shape(new_slide, int(slot_id))
            p = (plan_dir / img).resolve() if not Path(img).is_absolute() else Path(img)
            replace_picture(new_slide, shape, p)
        notes = entry.get("notes")
        if notes:
            new_slide.notes_slide.notes_text_frame.text = str(notes)
    # 删除原模板页(始终位于最前面的 n_template 页)
    for _ in range(n_template):
        delete_slide(prs, 0)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("plan")
    ap.add_argument("--check", action="store_true", help="只校验不生成")
    ap.add_argument("--allow-residual", action="store_true",
                    help="把「未填槽位将残留模板文案」从错误降级为警告")
    ap.add_argument("--scaffold", action="store_true",
                    help="生成 plan 骨架(需 --template 与 --pages)")
    ap.add_argument("--template", help="--scaffold 用: 模板 pptx 路径")
    ap.add_argument("--pages", help="--scaffold 用: 逗号分隔的模板页码序列, 可重复, "
                                    "如 1,2,4,4,4,6")
    ap.add_argument("--deck-out", help="--scaffold 用: 写入骨架的成品输出路径")
    args = ap.parse_args()

    plan_path = Path(args.plan).resolve()

    if args.scaffold:
        if not args.template or not args.pages:
            print("错误: --scaffold 需要 --template 和 --pages", file=sys.stderr)
            sys.exit(2)
        tpl = Path(args.template).resolve()
        if not tpl.is_file():
            print(f"错误: 找不到模板 {tpl}", file=sys.stderr)
            sys.exit(2)
        try:
            pages = [int(p) for p in args.pages.replace("，", ",").split(",") if p.strip()]
        except ValueError:
            print(f"错误: --pages 须为逗号分隔的整数: {args.pages}", file=sys.stderr)
            sys.exit(2)
        scaffold(plan_path, tpl, pages, args.deck_out)
        return

    if not plan_path.is_file():
        print(f"错误: 找不到 plan 文件 {plan_path}", file=sys.stderr)
        sys.exit(2)
    plan = json.loads(plan_path.read_text(encoding="utf-8-sig"))
    plan_dir = plan_path.parent

    tpl = Path(plan["template"])
    if not tpl.is_absolute():
        tpl = (plan_dir / tpl).resolve()
    if not tpl.is_file():
        print(f"错误: 找不到模板 {tpl}", file=sys.stderr)
        sys.exit(2)

    prs = Presentation(str(tpl))
    errors, warnings = validate(plan, plan_dir, prs, args.allow_residual)
    for w in warnings:
        print(f"W: {w}")
    for e in errors:
        print(f"E: {e}")
    if errors:
        print(f"\n校验失败: {len(errors)} 个错误, 未生成文件。", file=sys.stderr)
        sys.exit(1)
    if args.check:
        print(f"\n校验通过({len(warnings)} 个警告)。")
        return

    build(plan, plan_dir, prs)

    out = Path(plan["output"])
    if not out.is_absolute():
        out = (plan_dir / out).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"\n已生成 {len(plan['slides'])} 页: {out}")
    if warnings:
        print(f"注意: 有 {len(warnings)} 个警告(见上), 相关页必须重点复核。")
    print("QA: 对成品跑 snapshot.py 逐页目检; 截图不可用时, 以上方警告清单逐条核对"
          "文案长度并在交付时说明未做视觉复核。")


if __name__ == "__main__":
    main()
