# -*- coding: utf-8 -*-
"""pptxkit — ppt-template-fill 的共享核心库。

职责:
- 递归遍历幻灯片形状(含组合内嵌套形状), 识别可填充槽位
- 槽位 ID 规则: 文本/图片槽 = str(shape_id); 表格单元格 = "{shape_id}!r{row}c{col}" (1-based)
- 保留字符/段落格式的文本替换(在 lxml 层操作 run)
- 整页克隆(深拷贝 spTree + 关系重映射), 用于"以模板页为版式"的生成模式
- 图片槽替换(换 blip 关系 + srcRect 等比裁切, 不重编码图片)
"""

import copy
import re

from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# ---------------------------------------------------------------- 形状遍历

def iter_shapes(shapes, _path=()):
    """深度优先遍历, 进入组合形状。yield (shape, path_names)。"""
    for sh in shapes:
        path = _path + (sh.name or f"shape{sh.shape_id}",)
        yield sh, path
        if sh.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from iter_shapes(sh.shapes, path)


def find_shape(slide, shape_id):
    for sh, _ in iter_shapes(slide.shapes):
        if sh.shape_id == shape_id:
            return sh
    return None


# ---------------------------------------------------------------- 槽位识别

_DECORATIVE_RE = re.compile(
    r"^\s*(?:\d{1,3}|LOGO|logo|[·•\-–—|/\\.,:;'\"“”‘’()（）\s]+)\s*$"
)

def is_decorative_text(text):
    """页码、LOGO 字样、纯符号等装饰性文本, 默认不作为内容槽。"""
    t = text.strip()
    if not t:
        return False  # 空占位符另行处理
    if len(t) <= 1:
        return True
    return bool(_DECORATIVE_RE.match(t))


def get_text(shape):
    """形状可见文本, 段落以 \\n 连接。"""
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def effective_font_pt(shape):
    """尽力取首个显式字号(pt); 取不到返回 None(继承自版式/母版)。"""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size.pt
        if para.font.size is not None:
            return para.font.size.pt
    return None


def estimate_budget(shape, sample_len, font_pt):
    """估算文本槽的字数预算(按 CJK 全角字宽粗估)。返回 (max_chars, max_lines)。"""
    try:
        w, h = int(shape.width), int(shape.height)
    except (TypeError, ValueError):
        w = h = None
    if font_pt and w and h:
        glyph_emu = font_pt * 12700          # 1pt = 12700 EMU; CJK 字宽≈字号
        chars_per_line = max(1, int(w / glyph_emu * 0.96))
        max_lines = max(1, int(h / (glyph_emu * 1.35)))
        est = chars_per_line * max_lines
        # 样本长度是最可靠的下限参照
        return max(est, sample_len), max_lines
    if sample_len:
        return max(4, round(sample_len * 1.2)), None
    return None, None


_ROLE_PATTERNS = [
    ("toc", re.compile(r"目\s*录|contents|agenda|议程|CONTENTS", re.I)),
    ("ending", re.compile(r"感谢|谢谢|thank|再见|敬请|恳请.{0,4}指正", re.I)),
    ("section", re.compile(r"^第?\s*[0-9一二三四五六七八九十]+\s*(章|部分|节)|^PART\b", re.I | re.M)),
]

def guess_slide_role(slide, index, total):
    """粗分类: cover / toc / section / ending / content。仅作提示, 由 agent 结合截图定夺。"""
    layout_name = (slide.slide_layout.name or "").lower()
    all_text = "\n".join(get_text(sh) for sh, _ in iter_shapes(slide.shapes))
    for role, pat in _ROLE_PATTERNS:
        if pat.search(all_text):
            return role
    if "封面" in layout_name or "cover" in layout_name or (
        "title" in layout_name and index == 1
    ):
        return "cover"
    if "目录" in layout_name:
        return "toc"
    if "章节" in layout_name or "过渡" in layout_name or "section" in layout_name:
        return "section"
    if "结尾" in layout_name or "感谢" in layout_name or "end" in layout_name:
        return "ending"
    if index == 1:
        return "cover"
    if index == total:
        return "ending"
    return "content"


def collect_slots(slide):
    """收集一页内所有槽位。返回 dict: text_slots / image_slots / table_slots / other。"""
    text_slots, image_slots, table_slots, other = [], [], [], []
    for sh, path in iter_shapes(slide.shapes):
        if sh.shape_type == 6:  # 组合本身不是槽
            continue
        pos = {}
        try:
            pos = {
                "box_cm": [round(Emu(sh.width).cm, 1), round(Emu(sh.height).cm, 1)],
                "pos_cm": [round(Emu(sh.left).cm, 1), round(Emu(sh.top).cm, 1)],
            }
        except (TypeError, ValueError):
            pass
        if getattr(sh, "has_table", False) and sh.has_table:
            tbl = sh.table
            cells = {}
            for r, row in enumerate(tbl.rows, 1):
                for c in range(len(tbl.columns)):
                    cells[f"r{r}c{c + 1}"] = tbl.cell(r - 1, c).text
            table_slots.append({
                "id": str(sh.shape_id), "name": sh.name,
                "rows": len(tbl.rows), "cols": len(tbl.columns),
                "cells": cells, **pos,
            })
            continue
        if sh.shape_type == 13:  # PICTURE
            aspect = None
            try:
                aspect = round(int(sh.width) / int(sh.height), 2)
            except (TypeError, ValueError, ZeroDivisionError):
                pass
            image_slots.append({
                "id": str(sh.shape_id), "name": sh.name, "aspect": aspect, **pos,
            })
            continue
        if getattr(sh, "has_chart", False) and sh.has_chart:
            other.append({"id": str(sh.shape_id), "name": sh.name, "kind": "chart"})
            continue
        if sh.has_text_frame:
            sample = get_text(sh)
            is_ph = sh.is_placeholder
            if not sample.strip() and not is_ph:
                continue  # 无字的普通文本框视为装饰
            # 占位符可能不带自身尺寸/字号/提示文字, 需回退到版式同 idx 占位符
            layout_prompt = None
            layout_ph = None
            if is_ph:
                try:
                    idx = sh.placeholder_format.idx
                    for lph in slide.slide_layout.placeholders:
                        if lph.placeholder_format.idx == idx:
                            layout_ph = lph
                            break
                except (AttributeError, KeyError):
                    pass
            if layout_ph is not None and not sample.strip():
                layout_prompt = get_text(layout_ph) or layout_ph.name
            font_pt = effective_font_pt(sh)
            if font_pt is None and layout_ph is not None:
                font_pt = effective_font_pt(layout_ph)
            if font_pt is None and is_ph:
                # 字号继承自母版/主题时拿不到显式值, 按占位符类型给保守缺省
                try:
                    ph_type = str(sh.placeholder_format.type or "")
                except (AttributeError, KeyError):
                    ph_type = ""
                font_pt = 32.0 if "TITLE" in ph_type.upper() else 18.0
            size_ref = sh
            try:
                int(sh.width), int(sh.height)
            except (TypeError, ValueError):
                if layout_ph is not None:
                    size_ref = layout_ph
            sample_len = len(sample.replace("\n", ""))
            max_chars, max_lines = estimate_budget(size_ref, sample_len, font_pt)
            slot = {
                "id": str(sh.shape_id), "name": sh.name,
                "kind": "placeholder" if is_ph else "textbox",
                "sample": sample, "sample_len": sample_len,
                "paragraphs": len(sh.text_frame.paragraphs),
                "decorative": is_decorative_text(sample),
                **pos,
            }
            if layout_prompt:
                slot["layout_prompt"] = layout_prompt
            if font_pt:
                slot["font_pt"] = font_pt
            if max_chars:
                slot["max_chars"] = max_chars
            if max_lines:
                slot["max_lines"] = max_lines
            if len(path) > 1:
                slot["group_path"] = " > ".join(path[:-1])
            text_slots.append(slot)
    return {
        "text_slots": text_slots, "image_slots": image_slots,
        "table_slots": table_slots, "other": other,
    }


# ---------------------------------------------------------------- 文本替换

def _para_runs(p):
    return p.findall(qn("a:r"))


def _clear_para_content(p):
    """删除段落里所有 run / 换行 / 字段, 保留 pPr 与 endParaRPr。"""
    for tag in ("a:r", "a:br", "a:fld"):
        for el in p.findall(qn(tag)):
            p.remove(el)


def _proto_rPr(p):
    """为无 run 段落找一个 rPr 原型: 优先 endParaRPr, 其次 pPr/defRPr。"""
    end = p.find(qn("a:endParaRPr"))
    if end is not None:
        rpr = copy.deepcopy(end)
        rpr.tag = qn("a:rPr")
        return rpr
    return None


def set_paragraph_text(p, text):
    """替换单个段落文本, 保留首 run 的字符格式。"""
    runs = _para_runs(p)
    if runs:
        first = copy.deepcopy(runs[0])
        t = first.find(qn("a:t"))
        if t is None:
            t = first.makeelement(qn("a:t"), {})
            first.append(t)
        t.text = text
        _clear_para_content(p)
        end = p.find(qn("a:endParaRPr"))
        if end is not None:
            end.addprevious(first)
        else:
            p.append(first)
    else:
        _clear_para_content(p)
        r = p.makeelement(qn("a:r"), {})
        rpr = _proto_rPr(p)
        if rpr is not None:
            r.append(rpr)
        t = r.makeelement(qn("a:t"), {})
        t.text = text
        r.append(t)
        end = p.find(qn("a:endParaRPr"))
        if end is not None:
            end.addprevious(r)
        else:
            p.append(r)


def set_txbody_text(txBody, new_text):
    """替换整个 txBody 文本。new_text 以 \\n 分段; 多余模板段落删除, 不足则克隆末段。"""
    paras = txBody.findall(qn("a:p"))
    new_paras = str(new_text).split("\n")
    proto = paras[-1] if paras else None
    for i, ptext in enumerate(new_paras):
        if i < len(paras):
            p = paras[i]
        else:
            p = copy.deepcopy(proto)
            txBody.append(p)
        set_paragraph_text(p, ptext)
    for p in paras[len(new_paras):]:
        txBody.remove(p)


def set_shape_text(shape, new_text):
    if not shape.has_text_frame:
        raise ValueError(f"形状 {shape.shape_id} ({shape.name}) 没有文本框")
    set_txbody_text(shape.text_frame._txBody, new_text)


def set_table_cell_text(shape, row, col, new_text):
    cell = shape.table.cell(row - 1, col - 1)
    set_txbody_text(cell.text_frame._txBody, new_text)


# ---------------------------------------------------------------- 图片替换

def replace_picture(slide, pic_shape, image_path):
    """替换图片槽: 新增图片 part, 改 r:embed, 并用 a:srcRect 等比居中裁切避免拉伸。"""
    image_part, rId = slide.part.get_or_add_image_part(str(image_path))
    blipFill = pic_shape._element.blipFill
    blip = blipFill.find(qn("a:blip"))
    if blip is None:
        raise ValueError(f"图片形状 {pic_shape.shape_id} 缺少 a:blip")
    blip.set(qn("r:embed"), rId)
    # 清掉旧裁切, 按新图与槽位宽高比重新计算居中裁切
    for tag in ("a:srcRect", "a:tile", "a:stretch"):
        for el in blipFill.findall(qn(tag)):
            blipFill.remove(el)
    try:
        img_w, img_h = image_part.image.size  # 像素
        frame_ratio = int(pic_shape.width) / int(pic_shape.height)
        img_ratio = img_w / img_h
        l = r = t = b = 0
        if img_ratio > frame_ratio:
            cut = (1 - frame_ratio / img_ratio) / 2
            l = r = int(round(cut * 100000))
        elif img_ratio < frame_ratio:
            cut = (1 - img_ratio / frame_ratio) / 2
            t = b = int(round(cut * 100000))
        if any((l, r, t, b)):
            src = blipFill.makeelement(qn("a:srcRect"), {})
            for k, v in (("l", l), ("t", t), ("r", r), ("b", b)):
                if v:
                    src.set(k, str(v))
            blip.addnext(src)
    except (TypeError, ValueError, ZeroDivisionError):
        pass
    stretch = blipFill.makeelement(qn("a:stretch"), {})
    stretch.append(blipFill.makeelement(qn("a:fillRect"), {}))
    blipFill.append(stretch)


# ---------------------------------------------------------------- 整页克隆

_SKIP_REL_TYPES = {RT.SLIDE_LAYOUT, RT.NOTES_SLIDE}

def duplicate_slide(prs, source_slide):
    """在演示文稿末尾克隆 source_slide(同版式), 返回新 slide。

    深拷贝 spTree 后, 把源页的关系(图片/媒体/超链接等)在新页上重建,
    并按新 rId 重写 XML 里的 r:embed / r:id / r:link 引用。
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    # add_slide 内部已访问过 new_slide.shapes(lazyproperty 缓存绑定当前 spTree 节点),
    # 因此必须原地替换 spTree 的子节点, 不能替换 spTree 节点本身。
    src_copy = copy.deepcopy(source_slide._element.cSld.spTree)
    dst_spTree = new_slide._element.cSld.spTree
    for child in list(dst_spTree):
        dst_spTree.remove(child)
    for child in list(src_copy):
        dst_spTree.append(child)
    new_spTree = dst_spTree

    rid_map = {}
    for rId, rel in list(source_slide.part.rels.items()):
        if rel.reltype in _SKIP_REL_TYPES:
            continue
        if rel.is_external:
            new_rId = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rId] = new_rId

    if rid_map:
        for el in new_spTree.iter():
            for attr, val in list(el.attrib.items()):
                if attr.startswith("{%s}" % R_NS) and val in rid_map:
                    el.set(attr, rid_map[val])
    return new_slide


def delete_slide(prs, index):
    """按 0-based 索引删除一页。"""
    sldIdLst = prs.slides._sldIdLst
    slide_id = list(sldIdLst)[index]
    rId = slide_id.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(slide_id)
