# -*- coding: utf-8 -*-
"""inspect_template.py — 解析 PPTX 模板, 输出每页可填充槽位的清单(manifest)。

用法:
    python inspect_template.py <template.pptx> [--out manifest.json] [--compact]

- 默认输出完整 JSON manifest 到 stdout(或 --out 文件)。
- --compact 输出每页一行的摘要, 供快速选页; 细节仍以完整 manifest 为准。
"""

import argparse
import io
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from pptxkit import collect_slots, guess_slide_role, safe_layout  # noqa: E402


def build_manifest(path):
    prs = Presentation(path)
    total = len(prs.slides)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slots = collect_slots(slide)
        fillable = [s for s in slots["text_slots"] if not s["decorative"]]
        layout = safe_layout(slide)
        slides.append({
            "index": i,
            "layout_name": layout.name if layout is not None else "(无法解析)",
            "role_guess": guess_slide_role(slide, i, total),
            "fillable_text_slots": len(fillable),
            "capacity_chars": sum(s["sample_len"] for s in fillable),
            **slots,
        })
    return {
        "template": str(Path(path).resolve()),
        "slide_size_cm": [
            round(Emu(prs.slide_width).cm, 2),
            round(Emu(prs.slide_height).cm, 2),
        ],
        "slide_count": total,
        "slides": slides,
    }


def print_compact(manifest):
    print(f"template: {manifest['template']}")
    print(f"slides: {manifest['slide_count']}  "
          f"size: {manifest['slide_size_cm'][0]}x{manifest['slide_size_cm'][1]}cm")
    for s in manifest["slides"]:
        texts = [t for t in s["text_slots"] if not t["decorative"]]
        preview = " | ".join(
            (t["sample"].replace("\n", "⏎")[:14] or "<空占位符>") for t in texts[:4]
        )
        extras = []
        if s["image_slots"]:
            extras.append(f"图片x{len(s['image_slots'])}")
        if s["table_slots"]:
            t0 = s["table_slots"][0]
            extras.append(f"表格{t0['rows']}x{t0['cols']}")
        if s["other"]:
            extras.append("图表/其他x%d" % len(s["other"]))
        extra_str = (" [" + ",".join(extras) + "]") if extras else ""
        print(f"  p{s['index']:02d} {s['role_guess']:<8} "
              f"文本槽x{s['fillable_text_slots']}{extra_str}  "
              f"版式'{s['layout_name']}'  {preview}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("template")
    ap.add_argument("--out", help="manifest JSON 输出路径")
    ap.add_argument("--compact", action="store_true", help="只打印每页一行摘要")
    args = ap.parse_args()

    if not Path(args.template).is_file():
        print(f"错误: 找不到模板文件 {args.template}", file=sys.stderr)
        sys.exit(2)

    manifest = build_manifest(args.template)
    if args.out:
        Path(args.out).parent.mkdir(parents=True, exist_ok=True)
        Path(args.out).write_text(
            json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        print(f"manifest 已写入: {args.out}")
    if args.compact:
        print_compact(manifest)
    elif not args.out:
        print(json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
